#!/usr/bin/env python3
"""
Document Reviewer Tool
Extracts and reviews content from Word documents and PowerPoint presentations
"""

import os
from docx import Document
from pptx import Presentation
import sys

def extract_docx_content(file_path):
    """Extract text content from a Word document"""
    try:
        doc = Document(file_path)
        content = []
        
        print(f"\n{'='*60}")
        print(f"DOCUMENT: {os.path.basename(file_path)}")
        print(f"{'='*60}")
        
        for i, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                content.append(paragraph.text)
                print(f"Paragraph {i+1}: {paragraph.text}")
        
        # Extract table content if any
        for i, table in enumerate(doc.tables):
            print(f"\nTable {i+1}:")
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                if any(row_data):  # Only print non-empty rows
                    print(f"  {' | '.join(row_data)}")
        
        return content
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return []

def extract_pptx_content(file_path):
    """Extract text content from a PowerPoint presentation"""
    try:
        prs = Presentation(file_path)
        content = []
        
        print(f"\n{'='*60}")
        print(f"PRESENTATION: {os.path.basename(file_path)}")
        print(f"{'='*60}")
        
        for i, slide in enumerate(prs.slides):
            print(f"\nSlide {i+1}:")
            slide_content = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)
                    print(f"  - {shape.text}")
            
            content.append(slide_content)
        
        return content
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return []

def main():
    review_dir = "/Users/VScode_Projects/review"
    
    if not os.path.exists(review_dir):
        print(f"Directory {review_dir} not found!")
        return
    
    print("Week One Course Materials Review")
    print("=" * 40)
    
    # Process all files in the review directory
    for filename in os.listdir(review_dir):
        file_path = os.path.join(review_dir, filename)
        
        if filename.endswith('.docx'):
            extract_docx_content(file_path)
        elif filename.endswith('.pptx'):
            extract_pptx_content(file_path)
    
    print(f"\n{'='*60}")
    print("Review completed!")
    print(f"{'='*60}")

if __name__ == "__main__":
    main()
